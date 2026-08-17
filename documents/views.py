# Modules externes
import os
import zipfile
from openpyxl import load_workbook
from docx import Document as DocxDocument
from pptx import Presentation

# Modules Django
from django.shortcuts import render, get_object_or_404, redirect
from django.contrib.auth.decorators import login_required
from django.contrib import messages
from django.conf import settings
from django.core.exceptions import ValidationError

from .models import (
    Document_mfb, Document_dpse_planification, Document_dpse_suivi_evaluation,
    Document_dpse_statistiques, Document_dpse_qualite, Document_dpse_archives
)

# Mapping entre document_type et le modèle correspondant
DOCUMENT_MODELS = {
    'mfb': Document_mfb,
    'dpse_planification': Document_dpse_planification,
    'dpse_suivi_evaluation': Document_dpse_suivi_evaluation,
    'dpse_statistiques': Document_dpse_statistiques,
    'dpse_qualite': Document_dpse_qualite,
    'dpse_archives': Document_dpse_archives,
}

# Liste des documents
def document_list(request, document_type):

    Document = DOCUMENT_MODELS.get(document_type)

    documents = Document.objects.all()
    return render(request, 'documents/document_list.html', {'documents': documents,'document_type':document_type})

# Détails d'un document avec gestion des commentaires
def document_detail(request, document_type, pk):

    Document = DOCUMENT_MODELS.get(document_type)
    
    # Récupérer le document en utilisant pk
    document = get_object_or_404(Document, pk=pk)

    # Récupérer le chemin du fichier
    file_path = document.file.path

    """
    Fonction qui extrait le contenu d'un fichier (Word, Excel, PDF, PowerPoint ou image)
    et le retourne sous forme de liste de blocs de contenu.
    """
    content_blocks = []  # Liste qui stockera le contenu extrait sous forme de blocs

    # 📄 Gestion des fichiers Word (.docx)
    if file_path.endswith('.docx'):  
        doc = DocxDocument(file_path)  # Ouverture du fichier Word
        for element in doc.element.body:  # Parcours des éléments du document
            if element.tag.endswith('p'):  # Si l'élément est un paragraphe
                para = element  
                text = ''.join([node.text for node in para.iter() if node.text])  # Extraction du texte du paragraphe
                if text.strip():  # Vérification que le texte n'est pas vide
                    content_blocks.append({'type': 'text', 'content': text})  # Ajout du texte dans la liste

            elif element.tag.endswith('tbl'):  # Si l'élément est un tableau
                table_html = '<table class="table table-striped">'  # Début de la structure HTML du tableau
                for row in element.iterfind('.//w:tr', namespaces=element.nsmap):  # Parcours des lignes du tableau
                    table_html += '<tr>'  # Début d'une ligne
                    for cell in row.iterfind('.//w:tc', namespaces=element.nsmap):  # Parcours des cellules
                        cell_text = ''.join(node.text for node in cell.iterfind('.//w:t', namespaces=element.nsmap) if node.text)  # Extraction du texte de la cellule
                        table_html += f'<td>{cell_text}</td>'  # Ajout du texte dans une cellule HTML
                    table_html += '</tr>'  # Fin de la ligne
                table_html += '</table>'  # Fin du tableau
                content_blocks.append({'type': 'table', 'content': table_html})  # Ajout du tableau à la liste

        # 📷 Extraction des images contenues dans le fichier .docx
        if zipfile.is_zipfile(file_path):  # Vérification si le fichier est un ZIP (format interne de .docx)
            with zipfile.ZipFile(file_path, 'r') as docx_zip:  # Ouverture du fichier .docx comme une archive ZIP
                for file in docx_zip.namelist():  # Parcours des fichiers dans l'archive
                    if file.startswith('word/media/') and file.endswith(('png', 'jpeg', 'jpg', 'gif')):  # Si le fichier est une image
                        image_filename = os.path.basename(file)  # Récupération du nom du fichier image
                        image_path = os.path.join("documents", "files", document_type, "images", image_filename)  # Définition du chemin de sauvegarde
                        with open(image_path, 'wb') as image_file:  # Ouverture du fichier en mode écriture binaire
                            image_file.write(docx_zip.read(file))  # Écriture du fichier image
                        image_url = os.path.join(settings.MEDIA_URL, document_type,'documents_images', image_filename)  # Définition de l'URL de l'image
                        content_blocks.append({'type': 'image', 'content': image_url})  # Ajout de l'image à la liste

    # 📊 Gestion des fichiers Excel (.xlsx)
    elif file_path.endswith('.xlsx'):
        wb = load_workbook(file_path)  # Chargement du fichier Excel
        for sheet in wb.sheetnames:  # Parcours des feuilles du fichier Excel
            worksheet = wb[sheet]  
            table_html = '<table class="table table-striped">'  # Début d'un tableau HTML
            for row in worksheet.iter_rows(values_only=True):  # Parcours des lignes de la feuille Excel
                table_html += '<tr>'  # Début d'une ligne
                for cell in row:  # Parcours des cellules
                    table_html += f'<td>{cell if cell else ""}</td>'  # Ajout de la cellule au tableau
                table_html += '</tr>'  # Fin de la ligne
            table_html += '</table>'  # Fin du tableau
            content_blocks.append({'type': 'table', 'content': table_html})  # Ajout du tableau à la liste

    # 📜 Gestion des fichiers PDF (.pdf)
    elif file_path.endswith('.pdf'):
        content_blocks.append({'type': 'pdf', 'content': document.file.url})  # Ajout du fichier PDF (affichage en iframe ou lien)

    # 📽️ Gestion des fichiers PowerPoint (.pptx)
    elif file_path.endswith('.pptx'):
        prs = Presentation(file_path)  # Ouverture du fichier PowerPoint
        ppt_content = []  # Liste pour stocker le texte des diapositives
        for slide_number, slide in enumerate(prs.slides):  # Parcours des diapositives
            slide_text = []  # Liste pour stocker le texte d'une diapositive
            for shape in slide.shapes:  # Parcours des éléments de la diapositive
                if hasattr(shape, "text") and shape.text.strip():  # Si l'élément contient du texte
                    slide_text.append(shape.text.strip())  # Ajout du texte dans la liste

                # 📷 Extraction des images des diapositives
                if hasattr(shape, "image"):  # Si l'élément est une image
                    image_filename = f'slide_{slide_number + 1}.jpg'  # Nom du fichier image
                    image_path = os.path.join("documents", "files", document_type, "images", image_filename)  # Chemin de sauvegarde
                    with open(image_path, 'wb') as img_file:  # Ouverture du fichier en mode écriture binaire
                        img_file.write(shape.image.blob)  # Écriture de l'image
                    image_url = os.path.join(settings.MEDIA_URL, document_type,'documents_images', image_filename)  # Définition de l'URL de l'image
                    content_blocks.append({'type': 'image', 'content': image_url})  # Ajout de l'image à la liste

            if slide_text:  # Si du texte a été extrait de la diapositive
                ppt_content.append("\n".join(slide_text))  # Ajout du texte de la diapositive

        if ppt_content:  # Si des diapositives contiennent du texte
            content_blocks.append({'type': 'text', 'content': "\n\n".join(ppt_content)})  # Ajout du texte des diapositives

    # 🖼️ Gestion des fichiers image (.jpg, .png, .gif)
    elif file_path.endswith(('.gif', '.jpeg', '.jpg', '.png')):
        content_blocks.append({'type': 'image', 'content': document.file.url})  # Ajout de l'image à la liste

    # 📁 Gestion des autres fichiers (non pris en charge)
    else:
        content_blocks.append({'type': 'file', 'content': document.file.url})  # Ajout du fichier en tant que lien

    return render(request, 'documents/document_detail.html', {
        'document': document,
        'content_blocks': content_blocks,
        'document_type':document_type
    })

# Modifications pour la vue add_document
def add_document(request, document_type):
    Document = DOCUMENT_MODELS.get(document_type)
    
    if request.method == 'POST':
        name = request.POST.get('name')
        structure = request.POST.get('structure')
        description = request.POST.get('description')
        pans = request.POST.get('pans')
        results = request.POST.get('results')
        file = request.FILES.get('file')
        cover = request.FILES.get('cover')  # Nouveau champ

        if not file:
            return render(request, 'documents/add_document.html', {
                'error': "Veuillez télécharger un fichier."
            })

        document = Document(
            name=name,
            structure=structure,
            description=description,
            pans=pans,
            results=results,
            file=file,
            cover=cover,  # Ajout de la couverture
            uploaded_by=request.user
        )

        try:
            document.full_clean()
            document.save()
            messages.success(request, "Document ajouté avec succès!")
            return redirect('document_list', document_type=document_type)
        except ValidationError as e:
            messages.error(request, f"Erreur de validation : {e}")

    return render(request, 'documents/add_document.html', {'document_type': document_type})

# Modifications pour la vue edit_document
@login_required
def edit_document(request, document_type, pk):
    Document = DOCUMENT_MODELS.get(document_type)
    document = get_object_or_404(Document, pk=pk)

    if request.method == 'POST':
        document.name = request.POST.get('name', document.name)
        document.structure = request.POST.get('structure', document.structure)
        document.description = request.POST.get('description', document.description)
        document.pans = request.POST.get('pans', document.pans)
        document.results = request.POST.get('results', document.results)

        # Gestion du fichier principal
        if 'file' in request.FILES:
            if document.file:
                os.remove(document.file.path)
            document.file = request.FILES['file']

        # Gestion de la couverture
        if 'cover' in request.FILES:
            if document.cover:
                os.remove(document.cover.path)
            document.cover = request.FILES['cover']

        try:
            document.full_clean()
            document.save()
            messages.success(request, "Document modifié avec succès!")
            return redirect('document_detail', document_type=document_type, pk=document.pk)
        except ValidationError as e:
            messages.error(request, f"Erreur de validation : {e}")

    return render(request, 'documents/edit_document.html', {
        'document': document,
        'document_type': document_type
    })

# Modifications pour la vue delete_document
@login_required
def delete_document(request, document_type, pk):
    Document = DOCUMENT_MODELS.get(document_type)
    document = get_object_or_404(Document, pk=pk)

    if request.method == 'POST':
        # Suppression des fichiers associés
        if document.file:
            if os.path.isfile(document.file.path):
                os.remove(document.file.path)
        if document.cover:
            if os.path.isfile(document.cover.path):
                os.remove(document.cover.path)
        
        document.delete()
        messages.success(request, "Document supprimé avec succès!")
        return redirect('document_list', document_type=document_type)

    return render(request, 'documents/delete_confirm.html', {
        'document': document,
        'document_type': document_type
    })